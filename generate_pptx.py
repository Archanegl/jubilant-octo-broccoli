#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generate_pptx.py

将仓库根目录中名为 presentation.pptx 的文本草稿（Markdown 风格）解析为真实的 PowerPoint .pptx 文件。

用法（在本地）：
1. pip install python-pptx
2. python generate_pptx.py                 # 使用默认输入/输出
   或
   python generate_pptx.py -i draft.txt -o output/presentation_real.pptx

生成的文件： 默认为 output/presentation_real.pptx

脚本说明：
- 解析包含多个“幻灯片 N — 标题”块的文本文件。每个块中的列表项被视为幻灯片要点。
- 以“备注：”开头的行会被当作该幻灯片的备注（插入到幻灯片备注区），支持多行备注。
- 在每张幻灯片的内容区底部加入一行占位文字“[示意图占位 - 在此插入图片，建议来源：示意图/光纤/雷达 图片]”。

注意事项：
- 输入草稿是纯文本（UTF-8）文件。如果你使用的文件扩展名是 .pptx，但文件实际上是二进制 PowerPoint 文件，脚本会报错并提示你改用文本草稿文件名（例如 presentation.md 或 presentation.txt）。
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import re
import os
import argparse
import sys

# 默认输入（保留原脚本的默认文件名以兼容现有仓库）
DEFAULT_INPUT = 'presentation.pptx'
DEFAULT_OUTPUT_DIR = 'output'
DEFAULT_OUTPUT_FILE = os.path.join(DEFAULT_OUTPUT_DIR, 'presentation_real.pptx')

# 标题行匹配：支持中文“幻灯片 1 — 标题”或英文 Slide 1 - Title 等
SLIDE_TITLE_RE = re.compile(r'^(?:幻灯片|Slide)\s*\d+\s*[—\-–]\s*(.+)$', flags=re.I)

NOTE_PREFIX = '备注：'

# 多种常见的项目符号前缀（注意顺序：较长的先检查）
BULLET_PREFIXES = (
    '- ',  # ASCII dash + space
    ' - ', # with leading space
    '– ',  # en dash
    '— ',  # em dash
    '• ',  # bullet
    '· ',  # middle dot
    '＊ ', # ideographic asterisk
    '－ ',  # fullwidth hyphen
)


def looks_like_binary_pptx(path):
    """简单检测文件是否为二进制 .pptx (zip) 而非纯文本。
    这里通过检查开头字节判断（ZIP 文件以 PK 开头）。"""
    try:
        with open(path, 'rb') as f:
            head = f.read(4)
            return head.startswith(b'PK')  # .pptx is a ZIP container => starts with PK
    except Exception:
        return False


def parse_draft(path):
    """解析草稿文本，返回 slides 列表：[{title: str, bullets: [str], note: str}]"""
    with open(path, 'r', encoding='utf-8') as f:
        lines = [ln.rstrip('\n') for ln in f]

    slides = []
    cur = None

    for ln in lines:
        if not ln:
            # 空行作为段落分隔，直接跳过（但保留在备注里会被合并）
            continue
        m = SLIDE_TITLE_RE.match(ln.strip())
        if m:
            # start new slide
            if cur:
                slides.append(cur)
            cur = {'title': m.group(1).strip(), 'bullets': [], 'note': ''}
            continue
        if cur is None:
            # 忽略标题前的其它内容
            continue

        stripped = ln.strip()

        # 备注行：以“备注：”开头
        if stripped.startswith(NOTE_PREFIX):
            note_text = stripped[len(NOTE_PREFIX):].strip()
            if cur['note']:
                cur['note'] += '\n' + note_text
            else:
                cur['note'] = note_text
            continue

        # 忽略分割线
        if stripped.startswith('---'):
            continue

        # bullet 前缀检测（支持多种破折号或项目符号）
        matched_bullet = False
        for pref in BULLET_PREFIXES:
            if stripped.startswith(pref):
                cur['bullets'].append(stripped[len(pref):].strip())
                matched_bullet = True
                break
        if matched_bullet:
            continue

        # 若不是以 - 开头，也可能是上一条要点的续行，或是新的要点（文本行）
        if stripped:
            if cur['bullets']:
                # 续到上一 bullet
                cur['bullets'][-1] += ' ' + stripped
            else:
                # 没有 bullet 时把它当作新的 bullet
                cur['bullets'].append(stripped)

    if cur:
        slides.append(cur)
    return slides


def _find_non_title_text_placeholder(slide):
    """在 slide 中查找非标题的第一个可写 text_frame 占位符或文本框（fallback）。"""
    # 先尝试常见的占位符索引 1（大多数模板正文位于此）
    try:
        ph = slide.placeholders[1]
        if hasattr(ph, 'text_frame'):
            return ph
    except Exception:
        pass

    # 遍历所有占位符，选择第一个非标题且可写的
    for ph in slide.placeholders:
        try:
            # placeholder_format 属性在占位符上
            pf = getattr(ph, 'placeholder_format', None)
            # 尝试排除标题占位符（名字里含 title 的对象通常是标题）
            if pf is not None:
                name = getattr(ph, 'name', '') or ''
                if 'title' in name.lower():
                    continue
            if hasattr(ph, 'text_frame'):
                return ph
        except Exception:
            continue

    # 最后回退到遍历所有 shapes，找一个带 text_frame 的 shape
    for shp in slide.shapes:
        if hasattr(shp, 'text_frame'):
            # 排除标题（标题通常在 slide.shapes.title）
            if getattr(shp, 'is_placeholder', False):
                name = getattr(shp, 'name', '') or ''
                if 'title' in name.lower():
                    continue
            return shp

    return None


def make_pptx(slides, out_path):
    prs = Presentation()
    # choose a title slide layout and a content layout (indices commonly 0 and 1)
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # cover slide from first slide content
    if slides:
        s0 = prs.slides.add_slide(title_layout)
        # 安全地写入 title
        try:
            s0.shapes.title.text = slides[0]['title']
        except Exception:
            # 如果模板没有 title 占位符，找第一个 text frame 写入
            ph = _find_non_title_text_placeholder(s0)
            if ph is not None:
                ph.text_frame.clear()
                ph.text_frame.text = slides[0]['title']

        if slides[0]['bullets']:
            # 尝试把前几条要点写到封面副标题/占位符里
            subtitle_text = '\n'.join(slides[0]['bullets'][:3])
            try:
                subtitle = s0.placeholders[1]
                subtitle.text = subtitle_text
            except Exception:
                ph = _find_non_title_text_placeholder(s0)
                if ph is not None:
                    tf = ph.text_frame
                    tf.clear()
                    tf.text = subtitle_text

        # notes （notes_slide 会在内部创建）
        try:
            notes = s0.notes_slide.notes_text_frame
            if slides[0]['note']:
                notes.text = slides[0]['note']
        except Exception:
            # 忽略无法写入备注的情况
            pass

    # subsequent slides
    for sl in slides[1:]:
        slide = prs.slides.add_slide(content_layout)
        # 标题
        try:
            slide.shapes.title.text = sl['title']
        except Exception:
            ph = _find_non_title_text_placeholder(slide)
            if ph is not None:
                ph.text_frame.clear()
                ph.text_frame.text = sl['title']

        # add bullets
        body_ph = None
        try:
            body_ph = slide.shapes.placeholders[1]
        except Exception:
            body_ph = _find_non_title_text_placeholder(slide)

        if body_ph is None:
            # 如果找不到正文占位符，创建一个文本框作为后备
            left = Inches(0.5)
            top = Inches(1.8)
            width = Inches(9)
            height = Inches(3.0)
            body_shape = slide.shapes.add_textbox(left, top, width, height)
            body = body_shape.text_frame
            body.clear()
        else:
            body = body_ph.text_frame
            # 清除模板中已有的示例文本
            try:
                body.clear()
            except Exception:
                # 有些 text_frame 没有 clear 方法
                body.text = ''

        for i, b in enumerate(sl['bullets']):
            if i == 0:
                p = body.paragraphs[0]
                p.text = b
            else:
                p = body.add_paragraph()
                p.text = b
            p.level = 0
            try:
                p.font.size = Pt(18)
            except Exception:
                pass

        # add image placeholder text box at bottom
        left = Inches(0.5)
        top = Inches(5.0)
        width = Inches(9)
        height = Inches(0.6)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = '[示意图占位 - 在此插入图片，建议来源：示意图/光纤/雷达 图片]'
        for para in tf.paragraphs:
            try:
                para.font.italic = True
                para.font.size = Pt(12)
            except Exception:
                pass
            try:
                para.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            except Exception:
                pass

        # notes
        try:
            notes = slide.notes_slide.notes_text_frame
            if sl['note']:
                notes.text = sl['note']
        except Exception:
            pass

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print('Saved PPTX to', out_path)


def main(argv=None):
    parser = argparse.ArgumentParser(description='从文本草稿生成 PPTX（需要 python-pptx）')
    parser.add_argument('-i', '--input', default=DEFAULT_INPUT,
                        help='输入草稿文本文件（默认: %(default)s）')
    parser.add_argument('-o', '--output', default=DEFAULT_OUTPUT_FILE,
                        help='输出 .pptx 文件路径（默认: %(default)s）')
    args = parser.parse_args(argv)

    input_path = args.input
    output_path = args.output

    if not os.path.exists(input_path):
        print(f'Error: input file not found: {input_path}')
        sys.exit(2)

    # 如果输入文件看起来像二进制 PPTX（PK zip header），提示用户改用文本草稿
    if looks_like_binary_pptx(input_path):
        print(f'Error: the input "{input_path}" looks like a binary .pptx file (ZIP).')
        print('This script expects a plain-text draft (e.g. presentation.md or presentation.txt).')
        sys.exit(3)

    slides = parse_draft(input_path)
    if not slides:
        print('No slides parsed from draft. Please check the input file format.')
        sys.exit(4)

    make_pptx(slides, output_path)


if __name__ == '__main__':
    main()
